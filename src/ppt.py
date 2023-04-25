import collections 
import collections.abc
import pptx
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN

# Create a PowerPoint presentation
presentation = pptx.Presentation()

# Slide 1: Title Slide
title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
title_slide.shapes.title.text = "SAGE: Your Sustainability Chatbot"
subtitle = title_slide.placeholders[1].text_frame
subtitle.text = "Helping people live more sustainably"

# Slide 2: Introduction
intro_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
intro_slide.shapes.title.text = "Introduction"
intro_points = [
    "User-friendly AI Chatbot",
    "Quick and accurate responses",
    "Tailored advice on sustainable living",
]
for idx, point in enumerate(intro_points):
    intro_slide.placeholders[1].text_frame.add_paragraph().text = point

# Slide 3: Problem Statement
problem_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
problem_slide.shapes.title.text = "Problem Statement"
problem_points = [
    "Distinguish fad from fact",
    "Address environmental sustainability issues",
    "Offer credible, personalized solutions",
]
for idx, point in enumerate(problem_points):
    problem_slide.placeholders[1].text_frame.add_paragraph().text = point

# Slide 4: Problem Analysis
analysis_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
analysis_slide.shapes.title.text = "Problem Analysis"
analysis_points = [
    "Hybrid AI chatbot (machine learning and rule-based)",
    "Utilize NLP and expert systems",
    "Train on large datasets and user feedback",
]
for idx, point in enumerate(analysis_points):
    analysis_slide.placeholders[1].text_frame.add_paragraph().text = point

# Slide 5: Evaluation
evaluation_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
evaluation_slide.shapes.title.text = "Evaluation"
evaluation_points = [
    "Accuracy of information",
    "Relevance and completeness of responses",
    "Metrics: precision, recall, F1 score",
]
for idx, point in enumerate(evaluation_points):
    evaluation_slide.placeholders[1].text_frame.add_paragraph().text = point

# ... (previous code)

# Slide 6: Model and Dataset
model_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
model_slide.shapes.title.text = "Model and Dataset"
model_points = [
    "DistilBertForSequenceClassification (Hugging Face Transformers)",
    "Custom ChatDataset class for data preprocessing",
    "DataLoader for efficient training",
    "Training hyperparameters: 3 epochs, batch size 16, learning rate 5e-5",
    "Model trained using CrossEntropyLoss and AdamW optimizer",
]
for idx, point in enumerate(model_points):
    model_slide.placeholders[1].text_frame.add_paragraph().text = point

# Slide 7: Implementation Details
implementation_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
implementation_slide.shapes.title.text = "Implementation Details"
implementation_points = [
    "Import required packages and pre-trained models",
    "Load intents data from JSON file",
    "Tokenize and preprocess data using DistilBertTokenizer",
    "Create and train custom ChatDataset",
    "Train model using DataLoader, device-aware training, and optimizer",
    "Save trained model state and tags",
]
for idx, point in enumerate(implementation_points):
    implementation_slide.placeholders[1].text_frame.add_paragraph().text = point


# Save the presentation
presentation.save("SAGE_Sustainability_Chatbot.pptx")
